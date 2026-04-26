#!/usr/bin/env python3
"""
NotebookLM スライド後処理スクリプト
- スピーカーノート挿入
- 複数バッチの結合
- Insertスライドのマージ（チャンネルイントロ、メンバーシップCTA等）

Usage:
  python postprocess_slides.py \
    --speaker-notes speaker_notes.md \
    --output slides_merged_all.pptx \
    slides_batch1.pptx slides_batch2.pptx

  Insertスライド付き:
  python postprocess_slides.py \
    --speaker-notes speaker_notes.md \
    --insert templates/Insert.pptx \
    --insert-after 3 \
    --output MyProject.pptx \
    slides_batch1_slide1-16.pptx

  各バッチのスライド番号範囲はファイル名から自動検出:
    slides_batch1_slide1-20.pptx → Slide 1-20
    slides_batch2_slide21-32.pptx → Slide 21-32
"""

import argparse
import os
import re
import sys
import tempfile
from io import BytesIO

from pptx import Presentation


def parse_speaker_notes(notes_path: str) -> dict[int, str]:
    """speaker_notes.md を解析して {slide_num: note_text} を返す"""
    with open(notes_path, "r", encoding="utf-8") as f:
        content = f.read()

    notes = {}
    parts = re.split(r"【Slide (\d+)｜[^】]+】", content)
    for i in range(1, len(parts) - 1, 2):
        slide_num = int(parts[i])
        note_text = parts[i + 1].strip().rstrip("---").strip()
        notes[slide_num] = note_text

    return notes


def detect_slide_range(filename: str) -> tuple[int, int] | None:
    """ファイル名から slide{start}-{end} を検出"""
    match = re.search(r"slide(\d+)-(\d+)", filename)
    if match:
        return int(match.group(1)), int(match.group(2))
    return None


def insert_speaker_notes(prs: Presentation, notes: dict[int, str], start_slide: int) -> int:
    """PPTXにスピーカーノートを挿入。挿入数を返す"""
    inserted = 0
    for i, slide in enumerate(prs.slides):
        slide_num = start_slide + i
        if slide_num in notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes[slide_num]
            inserted += 1
    return inserted


def merge_batches(base_path: str, additional_paths: list[str], output_path: str) -> int:
    """バッチ1をベースに、追加バッチのスライドを結合"""
    base = Presentation(base_path)
    total_added = 0

    for path in additional_paths:
        prs = Presentation(path)
        for slide in prs.slides:
            # 画像を抽出
            img_shape = None
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    img_shape = shape
                    break

            if img_shape is None:
                print(f"  Warning: no image found in slide, skipping", file=sys.stderr)
                continue

            img_blob = img_shape.image.blob
            img_content_type = img_shape.image.content_type
            ext = "png" if "png" in img_content_type else "jpg"

            # ブランクスライドを追加して画像を配置
            blank_layout = base.slide_layouts[6]
            new_slide = base.slides.add_slide(blank_layout)

            # デフォルトプレースホルダーを削除
            for ph in list(new_slide.placeholders):
                sp = ph._element
                sp.getparent().remove(sp)

            # 画像をフルスクリーンで追加
            with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                tmp.write(img_blob)
                tmp_path = tmp.name

            new_slide.shapes.add_picture(
                tmp_path, 0, 0, base.slide_width, base.slide_height
            )
            os.unlink(tmp_path)

            # ノートをコピー
            if slide.has_notes_slide:
                src_notes = slide.notes_slide.notes_text_frame.text
                if src_notes.strip():
                    new_slide.notes_slide.notes_text_frame.text = src_notes

            total_added += 1

    base.save(output_path)
    return total_added


def _add_single_image_slide(dst_prs, src_slide, notes_text=""):
    """NotebookLMスライド用（1枚の画像のみ）"""
    slide_layout = dst_prs.slide_layouts[6]
    new_slide = dst_prs.slides.add_slide(slide_layout)
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    for shape in src_slide.shapes:
        if shape.shape_type == 13:  # Picture
            stream = BytesIO(shape.image.blob)
            new_slide.shapes.add_picture(stream, shape.left, shape.top, shape.width, shape.height)
    if notes_text:
        new_slide.notes_slide.notes_text_frame.text = notes_text
    return new_slide


def _add_complex_slide(dst_prs, src_slide, notes_text=""):
    """Insertスライド用（複数画像+テキストボックス対応）"""
    slide_layout = dst_prs.slide_layouts[6]
    new_slide = dst_prs.slides.add_slide(slide_layout)
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    for shape in src_slide.shapes:
        if shape.shape_type == 13:  # Picture
            stream = BytesIO(shape.image.blob)
            new_slide.shapes.add_picture(stream, shape.left, shape.top, shape.width, shape.height)
        elif shape.has_text_frame:
            txBox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            tf = txBox.text_frame
            tf.word_wrap = shape.text_frame.word_wrap
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                if p_idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.alignment = para.alignment
                for run in para.runs:
                    r = p.add_run()
                    r.text = run.text
                    if run.font.size:
                        r.font.size = run.font.size
                    if run.font.bold is not None:
                        r.font.bold = run.font.bold
                    if run.font.italic is not None:
                        r.font.italic = run.font.italic
                    if run.font.name:
                        r.font.name = run.font.name
                    try:
                        if run.font.color and run.font.color.type is not None and run.font.color.rgb:
                            r.font.color.rgb = run.font.color.rgb
                    except AttributeError:
                        pass

    if notes_text:
        new_slide.notes_slide.notes_text_frame.text = notes_text
    return new_slide


def _get_slide_notes(prs, slide_idx):
    """スライドからスピーカーノートを取得"""
    slide = prs.slides[slide_idx]
    if slide.has_notes_slide:
        return slide.notes_slide.notes_text_frame.text
    return ""


def merge_with_insert(main_path: str, insert_path: str, insert_after: int, output_path: str):
    """メインPPTXにInsertスライドをマージして新しいPPTXを出力

    Args:
        main_path: メインPPTXのパス
        insert_path: Insert.pptxのパス
        insert_after: Insert Slide 1を挿入する位置（N枚目の後）
        output_path: 出力PPTXのパス
    """
    main_prs = Presentation(main_path)
    insert_prs = Presentation(insert_path)

    output_prs = Presentation()
    output_prs.slide_width = main_prs.slide_width
    output_prs.slide_height = main_prs.slide_height

    total_main = len(main_prs.slides)
    total_insert = len(insert_prs.slides)

    print(f"  Main: {total_main} slides, Insert: {total_insert} slides")
    print(f"  Insert Slide 1 → after slide {insert_after}")
    print(f"  Insert Slide 2-{total_insert} → end")

    # Main slides 1 ~ insert_after
    for i in range(min(insert_after, total_main)):
        notes = _get_slide_notes(main_prs, i)
        _add_single_image_slide(output_prs, main_prs.slides[i], notes)

    # Insert Slide 1
    if total_insert >= 1:
        notes = _get_slide_notes(insert_prs, 0)
        _add_complex_slide(output_prs, insert_prs.slides[0], notes)

    # Main slides (insert_after+1) ~ end
    for i in range(insert_after, total_main):
        notes = _get_slide_notes(main_prs, i)
        _add_single_image_slide(output_prs, main_prs.slides[i], notes)

    # Insert Slides 2+ at the end
    for i in range(1, total_insert):
        notes = _get_slide_notes(insert_prs, i)
        _add_complex_slide(output_prs, insert_prs.slides[i], notes)

    output_prs.save(output_path)
    return len(output_prs.slides)


def main():
    parser = argparse.ArgumentParser(description="NotebookLM スライド後処理")
    parser.add_argument("pptx_files", nargs="+", help="バッチPPTXファイル（順番通りに指定）")
    parser.add_argument("--speaker-notes", "-s", help="speaker_notes.md のパス")
    parser.add_argument("--insert", "-i", help="Insert.pptx のパス（定型スライド挿入）")
    parser.add_argument("--insert-after", type=int, default=3,
                        help="Insert Slide 1を何枚目の後に挿入するか（デフォルト: 3 = イントロ後）")
    parser.add_argument("--output", "-o", required=True, help="出力ファイルパス")
    args = parser.parse_args()

    # スピーカーノート解析
    notes = {}
    if args.speaker_notes:
        notes = parse_speaker_notes(args.speaker_notes)
        print(f"Parsed {len(notes)} speaker notes")

    # 各バッチにスピーカーノートを挿入
    for path in args.pptx_files:
        prs = Presentation(path)
        filename = os.path.basename(path)
        slide_range = detect_slide_range(filename)

        if slide_range and notes:
            start = slide_range[0]
            inserted = insert_speaker_notes(prs, notes, start)
            prs.save(path)
            print(f"  {filename}: {inserted} notes inserted (Slide {slide_range[0]}-{slide_range[1]})")
        elif notes:
            # ファイル名からrange検出できない場合は1始まりと仮定
            inserted = insert_speaker_notes(prs, notes, 1)
            prs.save(path)
            print(f"  {filename}: {inserted} notes inserted (assumed Slide 1-)")

    # 結合（2ファイル以上の場合）
    if len(args.pptx_files) == 1:
        import shutil

        intermediate = args.output if not args.insert else args.output + ".tmp.pptx"
        shutil.copy2(args.pptx_files[0], intermediate)
        prs = Presentation(intermediate)
        print(f"\nOutput: {len(prs.slides)} slides")
    else:
        intermediate = args.output if not args.insert else args.output + ".tmp.pptx"
        added = merge_batches(args.pptx_files[0], args.pptx_files[1:], intermediate)
        prs = Presentation(intermediate)
        print(f"\nMerged: {len(prs.slides)} slides ({added} added from batches 2+)")

    # Insertスライドのマージ
    if args.insert:
        if not os.path.exists(args.insert):
            print(f"  Error: Insert file not found: {args.insert}", file=sys.stderr)
            sys.exit(1)

        print(f"\nInsert merge:")
        total = merge_with_insert(intermediate, args.insert, args.insert_after, args.output)
        print(f"  Final: {total} slides → {args.output}")

        # 中間ファイルを削除
        if intermediate != args.output and os.path.exists(intermediate):
            os.unlink(intermediate)
    else:
        # Insertなしの場合、intermediateが既にoutput
        prs = Presentation(args.output)
        print(f"  → {args.output}")

    print(f"File size: {os.path.getsize(args.output) / 1024 / 1024:.1f} MB")


if __name__ == "__main__":
    main()
